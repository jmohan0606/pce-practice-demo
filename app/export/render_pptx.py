"""PPTX renderer (python-pptx) — one slide per section, title carrying the
transition, table styled to the mockup tokens (navy header, green/red signed
columns, parenthesised negatives), traceability footer textbox."""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from app.export.payload import (NAVY, NEGATIVE_RED, NO_DATA_TEXT,
                                POSITIVE_GREEN, fmt_display, is_negative)

_NAVY = RGBColor.from_string(NAVY.lstrip("#"))
_POS = RGBColor.from_string(POSITIVE_GREEN.lstrip("#"))
_NEG = RGBColor.from_string(NEGATIVE_RED.lstrip("#"))
_SLATE = RGBColor.from_string("5A6B7D")
_INK = RGBColor.from_string("1A2430")

_MAX_BODY_ROWS = 18  # what fits legibly on one slide; overflow is stated, not hidden


def render_pptx(payload: dict) -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25),
                                         Inches(12.3), Inches(0.9))
    frame = title_box.text_frame
    frame.text = payload["title"]
    frame.paragraphs[0].runs[0].font.size = Pt(24)
    frame.paragraphs[0].runs[0].font.bold = True
    frame.paragraphs[0].runs[0].font.color.rgb = _NAVY
    sub = frame.add_paragraph()
    sub.text = payload["subtitle"]
    sub.runs[0].font.size = Pt(12)
    sub.runs[0].font.color.rgb = _SLATE

    top = Inches(1.25)
    if payload.get("preamble"):
        pre_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(1.1))
        pre_frame = pre_box.text_frame
        pre_frame.word_wrap = True
        for i, line in enumerate(payload["preamble"][:4]):
            para = pre_frame.paragraphs[0] if i == 0 else pre_frame.add_paragraph()
            para.text = line
            para.runs[0].font.size = Pt(10)
            para.runs[0].font.color.rgb = _INK
        top = Inches(2.4)

    columns = payload["columns"]
    rows = payload["rows"]
    body = list(rows)
    overflow_note = ""
    if len(body) > _MAX_BODY_ROWS:
        overflow_note = (f"Showing first {_MAX_BODY_ROWS} of {len(body)} rows — "
                         f"use the XLSX/CSV export for the full table.")
        body = body[:_MAX_BODY_ROWS]
    if payload.get("totals") and rows:
        body.append({**payload["totals"], "_role": "total"})

    n_rows = 1 + max(len(body), 1)
    table_height = Emu(int(Inches(0.32)) * n_rows)
    shape = slide.shapes.add_table(n_rows, len(columns), Inches(0.5), top,
                                   Inches(12.3), table_height)
    table = shape.table

    def _set_cell(cell, text: str, *, bold: bool = False,
                  color: RGBColor = _INK, fill: RGBColor | None = None) -> None:
        cell.text = text
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = Pt(10)
        run.font.bold = bold
        run.font.color.rgb = color
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill

    for c, column in enumerate(columns):
        _set_cell(table.cell(0, c), column["label"], bold=True,
                  color=RGBColor.from_string("FFFFFF"), fill=_NAVY)

    if not rows:
        _set_cell(table.cell(1, 0), NO_DATA_TEXT)
        for c in range(1, len(columns)):
            _set_cell(table.cell(1, c), "")
    else:
        for r, row in enumerate(body, start=1):
            emphatic = row.get("_role") in ("subtotal", "total")
            for c, column in enumerate(columns):
                value = row.get(column["key"])
                color = _INK
                if column.get("signed") and column["type"] != "text" and value is not None:
                    color = _NEG if is_negative(value) else _POS
                _set_cell(table.cell(r, c), fmt_display(value, column["type"]),
                          bold=emphatic, color=color,
                          fill=RGBColor.from_string("EDF1F6") if emphatic else None)

    footer_lines = []
    if overflow_note:
        footer_lines.append(overflow_note)
    footer_lines += (payload.get("footnotes") or [])[:2]
    footer = payload["footer"]
    footer_lines.append(f"Source: {footer['source']} · Generated: "
                        f"{footer['generated_at']} · Rule set version: "
                        f"{footer['rule_set_version']}")
    foot_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.75),
                                        Inches(12.3), Inches(0.7))
    foot_frame = foot_box.text_frame
    foot_frame.word_wrap = True
    for i, line in enumerate(footer_lines):
        para = foot_frame.paragraphs[0] if i == 0 else foot_frame.add_paragraph()
        para.text = line
        para.runs[0].font.size = Pt(8)
        para.runs[0].font.color.rgb = _SLATE

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()
