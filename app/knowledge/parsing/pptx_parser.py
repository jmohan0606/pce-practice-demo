"""PPTX parsing via python-pptx (Round B rework, spec B2.1).

One heading block per slide (the slide title, level 1) followed by one
paragraph block with the slide's remaining text; slide tables become their own
table blocks. page_no is the 1-based slide number.
"""

from __future__ import annotations

from pathlib import Path

from app.knowledge.parsing.base import ParsedBlock, ParsedDocument, table_to_markdown


def parse_pptx(path: str | Path) -> ParsedDocument:
    from pptx import Presentation

    path = Path(path)
    presentation = Presentation(str(path))
    blocks: list[ParsedBlock] = []
    order = 0

    slides = list(presentation.slides)
    for slide_no, slide in enumerate(slides, start=1):
        title_shape = slide.shapes.title
        # python-pptx returns a fresh proxy per access — identify by shape_id.
        title_shape_id = title_shape.shape_id if title_shape is not None else None
        title = (title_shape.text.strip() if title_shape is not None else "")
        if title:
            blocks.append(ParsedBlock(
                block_type="heading", text=title, page_no=slide_no,
                heading_level=1, order=order,
            ))
            order += 1

        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                rows = [r for r in rows if any(c.strip() for c in r)]
                markdown = table_to_markdown(rows)
                if markdown:
                    blocks.append(ParsedBlock(
                        block_type="table", text=markdown, page_no=slide_no,
                        heading_level=0, order=order,
                    ))
                    order += 1
                continue
            if shape.has_text_frame and shape.shape_id != title_shape_id:
                text = shape.text.strip()
                if text:
                    body_parts.append(text)
        if body_parts:
            blocks.append(ParsedBlock(
                block_type="paragraph", text="\n".join(body_parts), page_no=slide_no,
                heading_level=0, order=order,
            ))
            order += 1

    return ParsedDocument(source_name=path.name, page_count=len(slides), blocks=blocks)
