#!/usr/bin/env python3
"""Round A1 task 6 check — exports generate, OPEN, and CONTAIN CONTENT.

Every file is read back with an independent reader: pypdf extracts text (a PDF
that renders blank is a failure), python-pptx reads table cells, the xlsx is
parsed from its XML to prove numeric cells hold RAW numbers, and the csv is
re-parsed. Then POST /api/export is exercised for every section × format.

Runs against an isolated runtime DB (PCE_RUNTIME_DB_DIR) so it never touches
data/runtime.
"""
from __future__ import annotations

import csv
import io
import os
import re
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree

os.environ.setdefault("PCE_RUNTIME_DB_DIR", tempfile.mkdtemp(prefix="pce_export_check_"))

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

OUT_DIR = Path(tempfile.mkdtemp(prefix="pce_export_files_"))
PARAMS = {
    "dashboard_table": {"from": "202604", "to": "202605", "view": "all"},
    "noncredited": {"month": "202605"},
    "exceptions": {"from": "202604", "to": "202605"},
    "insights": {"from": "202604", "to": "202605", "advisor": "all"},
}

CHECKS: list[tuple[str, bool, str]] = []


def check(name: str, ok: bool, detail: str) -> None:
    CHECKS.append((name, ok, detail))
    print(f"[{'PASS' if ok else 'FAIL'}] {name}: {detail}")


def read_back_pdf(path: Path) -> None:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    snippet = " ".join(text.split())[:160]
    footer_ok = "Rule set version" in text and "Generated:" in text
    check(f"pdf read-back {path.name}",
          len(reader.pages) >= 1 and len(text.strip()) > 40 and footer_ok,
          f"{len(reader.pages)} page(s), {len(text)} chars extracted; "
          f"footer present={footer_ok}; snippet: {snippet!r}")
    footer_snip = " ".join(text.split())
    idx = footer_snip.find("Rule set version")
    print(f"       pdf footer text: ...{footer_snip[max(0, idx - 60):idx + 40]!r}")


def read_back_pptx(path: Path) -> None:
    from pptx import Presentation

    deck = Presentation(str(path))
    tables = [shape.table for slide in deck.slides
              for shape in slide.shapes if shape.has_table]
    cell00 = tables[0].cell(0, 0).text if tables else ""
    cell10 = tables[0].cell(1, 0).text if tables and len(tables[0].rows) > 1 else ""
    texts = " | ".join(shape.text_frame.text for slide in deck.slides
                       for shape in slide.shapes if shape.has_text_frame)
    check(f"pptx read-back {path.name}",
          len(deck.slides) == 1 and bool(tables) and bool(cell00)
          and "Rule set version" in texts,
          f"{len(deck.slides)} slide, {len(tables)} table(s), header cell "
          f"{cell00!r}, first body cell {cell10!r}, footer present="
          f"{'Rule set version' in texts}")


def read_back_xlsx(path: Path) -> None:
    """Prove numeric cells are RAW numbers with number formats, using only the
    file's own XML (openpyxl is not installed in this environment)."""
    ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(path) as z:
        sheet = ElementTree.fromstring(z.read("xl/worksheets/sheet1.xml"))
        styles = z.read("xl/styles.xml").decode()
    numeric = []
    for cell in sheet.iter("{%s}c" % ns["m"]):
        if cell.get("t") in ("s", "str", "inlineStr"):
            continue
        v = cell.find("m:v", ns)
        if v is not None and v.text is not None:
            numeric.append((cell.get("r"), float(v.text)))
    fmts = re.findall(r'formatCode="([^"]+)"', styles)
    has_paren_fmt = any("(" in f and ")" in f for f in fmts)
    check(f"xlsx read-back {path.name}",
          len(numeric) >= 3 and has_paren_fmt,
          f"{len(numeric)} raw numeric cells (e.g. {numeric[:3]}); "
          f"parenthesised-negative number formats present={has_paren_fmt} "
          f"(formats: {fmts[:3]})")


def read_back_csv(path: Path) -> None:
    rows = list(csv.reader(io.StringIO(path.read_text("utf-8-sig"))))
    flat = ["|".join(r) for r in rows]
    footer_rows = [r for r in flat if r.startswith("# rule_set_version")
                   or r.startswith("# generated_at") or r.startswith("# source")]
    check(f"csv read-back {path.name}",
          len(rows) >= 6 and len(footer_rows) == 3,
          f"{len(rows)} parsed rows; footer rows: {footer_rows}")


READERS = {"pdf": read_back_pdf, "pptx": read_back_pptx,
           "xlsx": read_back_xlsx, "csv": read_back_csv}


def main() -> int:
    print(f"runtime db dir : {os.environ['PCE_RUNTIME_DB_DIR']}")
    print(f"output dir     : {OUT_DIR}\n")

    # Building the app seeds rule-set v0 (published) into the isolated DB, so
    # the footer's rule_set_version is real, not a fixture.
    from fastapi.testclient import TestClient

    from app.api.main import create_app
    from app.export.service import FORMATS, SECTIONS, export_file

    client = TestClient(create_app())

    print("== direct generation + independent read-back ==")
    plan = {"dashboard_table": list(FORMATS),
            "noncredited": ["csv", "pdf", "xlsx", "pptx"],
            "exceptions": ["csv", "pdf"],
            "insights": ["csv", "pdf"]}
    for section, formats in plan.items():
        for fmt in formats:
            content, media_type, filename = export_file(section, fmt, PARAMS[section])
            path = OUT_DIR / filename
            path.write_bytes(content)
            print(f"\n-- {section} / {fmt}: {filename} "
                  f"({len(content):,} bytes, {media_type})")
            # csv has no container overhead — an honest "no data" table is small
            check(f"{section}.{fmt} size",
                  len(content) > (300 if fmt == "csv" else 700),
                  f"{len(content):,} bytes")
            READERS[fmt](path)

    print("\n== POST /api/export — every section × format ==")
    for section in SECTIONS:
        for fmt in FORMATS:
            resp = client.post("/api/export", json={
                "section": section, "format": fmt, "params": PARAMS[section]})
            disposition = resp.headers.get("content-disposition", "")
            check(f"POST {section}/{fmt}",
                  resp.status_code == 200 and len(resp.content) > 400
                  and "attachment" in disposition,
                  f"status={resp.status_code} bytes={len(resp.content):,} "
                  f"type={resp.headers.get('content-type')} "
                  f"disposition={disposition}")

    print("\n== error handling ==")
    bad = client.post("/api/export", json={"section": "nope", "format": "pdf",
                                           "params": {}})
    check("unknown section -> 400", bad.status_code == 400,
          f"status={bad.status_code} detail={bad.json().get('detail')!r}")
    bad = client.post("/api/export", json={"section": "dashboard_table",
                                           "format": "docx", "params": PARAMS["dashboard_table"]})
    check("unknown format -> 400", bad.status_code == 400,
          f"status={bad.status_code} detail={bad.json().get('detail')!r}")
    bad = client.post("/api/export", json={"section": "dashboard_table",
                                           "format": "pdf",
                                           "params": {"from": "999999", "to": "202605"}})
    check("unknown month -> 400", bad.status_code == 400,
          f"status={bad.status_code} detail={bad.json().get('detail')!r}")

    failed = [c for c in CHECKS if not c[1]]
    print(f"\n{'=' * 60}\n{len(CHECKS) - len(failed)}/{len(CHECKS)} checks passed")
    for name, _, detail in failed:
        print(f"  FAILED: {name} — {detail}")
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main())
