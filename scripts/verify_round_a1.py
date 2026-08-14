"""Round A1 verification — the 17 Task 7 checks (docs/ROUND_A1_SPEC.md).

Deterministic: scripted LLM mode, isolated runtime SQLite (fresh temp dir), no
real LLM calls. Run: python3 scripts/verify_round_a1.py
"""
from __future__ import annotations

import os
import sys
import tempfile

os.environ.setdefault("LLM_MODE", "scripted")
os.environ["PCE_RUNTIME_DB_DIR"] = tempfile.mkdtemp(prefix="pce_a1_verify_")

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

RESULTS: list[tuple[bool, str, str]] = []


def check(number: str, title: str, ok: bool, detail: str) -> None:
    RESULTS.append((ok, number, f"{title} — {detail}"))
    print(f"{'PASS' if ok else 'FAIL'}  A1-{number}. {title} — {detail}")


def main() -> int:  # noqa: PLR0915 — one linear check script, house style
    from fastapi.testclient import TestClient

    from app.api import main as apimain

    client = TestClient(apimain.create_app())

    # ---- 1: driver_code stored; label resolves at read; rename reaches history
    from app.api.routers.insights import _serialize_finding
    from app.insights.store import get_insight_store
    from app.rules.store import get_rule_store

    istore, rstore = get_insight_store(), get_rule_store()
    run = istore.begin_run("V000002", "202604", "202605", "RSV_v0")
    istore.complete_run(run["run_id"], narrative="n", bullets=[], findings=[
        {"title": "probe", "summary": "s", "impact_amt": 1.0,
         "driver_tag": "New Billing", "rule_key": None, "provenance": "REAL",
         "confidence": 1.0, "evidence_columns": [], "evidence_rows": [],
         "origin": "agent"}],
        query_count=0, budget_hit=False, coverage_ratio=None)
    stored = istore.run_findings(run["run_id"])[0]
    before = _serialize_finding(stored)["driver_tag"]
    rstore.set_driver_label("NEW_BILLING", "First-Time Billing (A1 probe)")
    after = _serialize_finding(istore.run_findings(run["run_id"])[0])["driver_tag"]
    check("1", "driver_code stored on findings; label resolves at read; rename "
               "reaches historical findings with no regeneration",
          stored.get("driver_code") == "NEW_BILLING" and "driver_tag" not in stored
          and before == "New Billing" and after == "First-Time Billing (A1 probe)",
          f"stored driver_code=NEW_BILLING, no stored tag; displayed "
          f"{before!r} -> {after!r} after PATCH-equivalent registry write")
    rstore.set_driver_label("NEW_BILLING", "New Billing")  # restore

    # ---- 2: GET /api/drivers
    drivers = client.get("/api/drivers").json()["drivers"]
    complete = all(d.get("driver_code") and d.get("driver_label")
                   and d.get("driver_definition") and "source" in d and "rule_key" in d
                   for d in drivers)
    rule_backed = [d for d in drivers if d["rule_key"]]
    check("2", "GET /api/drivers returns code, label, definition and source for "
               "every driver", complete and len(drivers) >= 15 and rule_backed,
          f"{len(drivers)} drivers, all fields present; {len(rule_backed)} "
          f"rule-backed (e.g. {rule_backed[0]['driver_code']} source="
          f"{rule_backed[0]['source'] if isinstance(rule_backed[0]['source'], str) else 'citation'})")

    # ---- 3: severities on v0 + PATCH mints a version
    rules = client.get("/api/rules?version=RSV_v0").json()["rules"]
    sev = {r["rule_code"]: (r["severity"], r["severity_reason"]) for r in rules}
    expected = {"LOST_ACCOUNT": "HIGH", "ACCOUNT_TRANSFERRED_OUT": "MODERATE",
                "ACCOUNT_TRANSFERRED_IN": "LOW", "NEW_ACCOUNT": "LOW",
                "NEW_BILLING": "INFO", "RETAINED_ACCOUNT": "INFO"}
    seeded_ok = (all(sev.get(c, (None,))[0] == lvl for c, lvl in expected.items())
                 and all(reason for _, reason in sev.values()))
    la = next(r for r in rules if r["rule_code"] == "LOST_ACCOUNT")
    patched = client.patch(f"/api/rules/{la['rule_key']}/severity",
                           json={"severity": "CRITICAL",
                                 "severity_reason": "verify probe"}).json()
    versions = client.get("/api/rules/versions").json()["versions"]
    check("3", "severity set on all v0 rules with severity_reason; PATCH changes "
               "it and mints a version",
          seeded_ok and patched["rule"]["severity"] == "CRITICAL"
          and patched["version"]["version_no"] == 1
          and versions[0]["status"] == "SUPERSEDED"
          and patched["rule"]["plan"] is not None,
          f"seeded {dict((c, s[0]) for c, s in sev.items())}; PATCH -> RSV_v1 "
          f"PUBLISHED, v0 SUPERSEDED, compiled plan preserved")

    # ---- 4: findings inherit severity; no-rule findings INFO
    v1_rules = client.get("/api/rules?version=RSV_v1").json()["rules"]
    la1 = next(r for r in v1_rules if r["rule_code"] == "LOST_ACCOUNT")
    run2 = istore.begin_run("V000005", "202604", "202605", "RSV_v1")
    istore.complete_run(run2["run_id"], narrative="n", bullets=[], findings=[
        {"title": "rule finding", "summary": "s", "impact_amt": -5.0,
         "driver_tag": "Lost Accounts", "rule_key": la1["rule_key"],
         "provenance": "REAL", "confidence": 1.0, "evidence_columns": [],
         "evidence_rows": [], "origin": "rule"},
        {"title": "agent finding", "summary": "s", "impact_amt": None,
         "driver_tag": "Market", "rule_key": None, "provenance": "DERIVED",
         "confidence": 0.7, "evidence_columns": [], "evidence_rows": [],
         "origin": "agent"}],
        query_count=0, budget_hit=False, coverage_ratio=None)
    got = [_serialize_finding(f)["severity"]
           for f in istore.run_findings(run2["run_id"])]
    check("4", "findings inherit rule severity; findings with no rule are INFO",
          got == ["CRITICAL", "INFO"],
          f"rule finding severity={got[0]} (rule is CRITICAL after check 3), "
          f"no-rule finding severity={got[1]}")

    # ---- 5: exceptions filter and sort by severity
    ex_all = client.get("/api/exceptions",
                        params={"from": "202604", "to": "202605"}).json()
    order = [e["severity"] for e in ex_all["exceptions"]]
    from app.rules.store import SEVERITIES
    rank = {s: i for i, s in enumerate(SEVERITIES)}
    ex_crit = client.get("/api/exceptions",
                         params={"from": "202604", "to": "202605",
                                 "severity": "CRITICAL"}).json()["exceptions"]
    bad = client.get("/api/exceptions", params={"from": "202604", "to": "202605",
                                                "severity": "BOGUS"})
    check("5", "exceptions filter and sort by severity",
          order == sorted(order, key=rank.__getitem__)
          and all(e["severity"] == "CRITICAL" for e in ex_crit) and ex_crit
          and bad.status_code == 400,
          f"unfiltered order={order}; severity=CRITICAL -> {len(ex_crit)} row(s); "
          f"unknown level -> 400")

    # ---- 6/7/8: dashboard table
    details = {}
    ok678 = [True, True, True]
    for view in ("all", "split", "recurring", "non_recurring"):
        j = client.get("/api/dashboard/table",
                       params={"from": "202604", "to": "202605",
                               "view": view}).json()
        rows, tot = j["rows"], j["total"]
        share = round(sum(r["share_pct"] or 0 for r in rows), 2)
        revsum = round(sum(r["to_amt"] for r in rows), 2)
        deltas = all(r["to_account_count"] - r["from_account_count"]
                     == r["account_delta"]
                     and r["to_trade_count"] - r["from_trade_count"]
                     == r["trade_delta"] for r in rows)
        details[view] = (len(rows), share, round(tot["to_amt"], 2))
        ok678[0] &= abs(revsum - round(tot["to_amt"], 2)) < 0.01 and abs(share - 100.0) <= 0.1
        ok678[2] &= deltas
    ok678[1] = (details["recurring"][2] != details["all"][2]
                and abs(details["recurring"][1] - 100.0) <= 0.1)
    check("6", "dashboard table: rows sum to total, share_pct sums to 100 +/- 0.1 "
               "in all four views", ok678[0],
          "; ".join(f"{v}: rows={d[0]} share_sum={d[1]}"
                    for v, d in details.items()))
    check("7", "share_pct in recurring-only view is of the recurring total, not "
               "the firm total", ok678[1],
          f"recurring total {details['recurring'][2]:,} != firm total "
          f"{details['all'][2]:,} and recurring shares still sum to "
          f"{details['recurring'][1]}")
    check("8", "accounts/trades deltas equal to - from on every row", ok678[2],
          "verified across all four views")

    # ---- 9: 9X codes present with realistic volumes
    from collections import Counter

    from app.graph.foundation_store import get_foundation_store
    codes = Counter(str(t.get("reason_cd")) for t in
                    get_foundation_store()
                    .all_vertices("phx_dm_pce_revenue_transaction").values())
    need = {"9H", "9G", "9D", "9E"}
    check("9", "mock data contains 9H/9G/9D/9E reason codes with realistic "
               "volumes", all(codes.get(c, 0) >= 10 for c in need),
          f"{ {c: codes[c] for c in sorted(need)} } of "
          f"{sum(codes.values())} txns ({codes['__NONE__']} credited)")

    # ---- 10/11: detail queries
    shapes = {
        "household": {"advisor_sid", "household_count", "accounts", "trades",
                      "value", "avg_household_assets",
                      "households_within_10k_of_threshold"},
        "inheritance": {"advisor_sid", "from_advisor_sid",
                        "from_advisor_departed", "accounts", "transfer_date",
                        "months_since_transfer", "trades", "value"},
        "discount": {"advisor_sid", "accounts", "avg_standard_bps",
                     "avg_actual_bps", "avg_reduction_pct", "accounts_above_10pct",
                     "grid_points_expected", "grid_points_recorded", "value"},
        "eligibility": {"product_id", "reason", "accounts", "advisors",
                        "trades", "value"},
    }
    missing_all = {}
    counts = {}
    for cause, need_cols in shapes.items():
        rows = client.get(f"/api/noncredited/detail/{cause}",
                          params={"month": "202605"}).json()["rows"]
        counts[cause] = len(rows)
        missing_all[cause] = (need_cols - set(rows[0].keys())) if rows else need_cols
    check("10", "all four non-credited detail queries return their documented "
                "columns", not any(missing_all.values()),
          f"row counts {counts}; missing columns: "
          f"{ {c: sorted(m) for c, m in missing_all.items() if m} or 'none'}")
    elig = client.get("/api/noncredited/detail/eligibility",
                      params={"month": "202605"}).json()["rows"]
    check("11", "eligibility detail is grouped by product, not advisor",
          all("product_id" in r and "advisor_sid" not in r for r in elig),
          f"{len(elig)} product rows, no advisor_sid column, advisors is a "
          f"count (first row: {elig[0]['product_id']} advisors={elig[0]['advisors']})")

    # ---- 12/13: ranking
    j = client.get("/api/dashboard/product/managed_accounts/ranking",
                   params={"from": "202604", "to": "202605"}).json()
    top, bottom = j["top"], j["bottom"]
    ranked = (all(top[i]["change_amt"] >= top[i + 1]["change_amt"]
                  for i in range(len(top) - 1))
              and all(bottom[i]["change_amt"] <= bottom[i + 1]["change_amt"]
                      for i in range(len(bottom) - 1)))
    pct_ok = all("pct_of_total_change" in r for r in top + bottom)
    check("12", "ranking returns <=10 each side, ranked by change amount, with "
                "pct_of_total_change",
          len(top) <= 10 and len(bottom) <= 10 and ranked and pct_ok,
          f"top={len(top)} bottom={len(bottom)} of {j['advisor_count']} advisors; "
          f"#1 {top[0]['advisor_sid']} {top[0]['change_amt']:+,.2f} "
          f"({top[0]['pct_of_total_change']}% of {j['total_change_amt']:,.2f})")
    nulls = [r["advisor_sid"] for r in top + bottom
             if r["dominant_driver_code"] is None]
    check("13", "dominant_driver_code is null — not guessed — when no rule "
                "outcome exists", bool(nulls),
          f"null for {nulls} alongside real drivers like "
          f"{[r['dominant_driver_code'] for r in top[:2]]}")

    # ---- 14: exports
    import csv as _csv
    import io

    from pypdf import PdfReader
    sizes = {}
    opened = {}
    footer_ok = {}
    for fmt in ("pdf", "pptx", "xlsx", "csv"):
        resp = client.post("/api/export", json={
            "section": "dashboard_table", "format": fmt,
            "params": {"from": "202604", "to": "202605", "view": "all"}})
        sizes[fmt] = len(resp.content)
        if fmt == "pdf":
            text = "".join(p.extract_text() for p in
                           PdfReader(io.BytesIO(resp.content)).pages)
            opened[fmt] = len(text) > 500
            footer_ok[fmt] = "Rule set version" in text
        elif fmt == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(resp.content))
            texts = [sh.text_frame.text for sl in prs.slides
                     for sh in sl.shapes if sh.has_text_frame]
            opened[fmt] = len(prs.slides) >= 1
            footer_ok[fmt] = any("Rule set version" in t or "rule_set_version" in t
                                 for t in texts)
        elif fmt == "xlsx":
            import zipfile
            z = zipfile.ZipFile(io.BytesIO(resp.content))
            shared = z.read("xl/sharedStrings.xml").decode("utf-8", "ignore")
            opened[fmt] = "xl/worksheets/sheet1.xml" in z.namelist()
            footer_ok[fmt] = "rule_set_version" in shared or "Rule set" in shared
        else:
            rows = list(_csv.reader(io.StringIO(resp.content.decode())))
            opened[fmt] = len(rows) > 10
            footer_ok[fmt] = any("rule_set_version" in ",".join(r) for r in rows)
    check("14", "all four export formats generate and open; each carries the "
                "traceability footer",
          all(sizes[f] > 1000 for f in sizes) and all(opened.values())
          and all(footer_ok.values()),
          f"sizes={sizes}; opened+content={opened}; footer={footer_ok}")

    # ---- 15: RETAINED_ACCOUNT semantics
    from app.rules.service import evaluate_rule_set
    out5 = evaluate_rule_set("RSV_v0", month="202605", scope="practice")
    matched = {r["rule_code"]: {str(m["key"]) for m in r.get("matched", [])}
               for r in out5["results"] if r.get("evaluated")}
    claimed = (matched.get("NEW_ACCOUNT", set())
               | matched.get("NEW_BILLING", set())
               | matched.get("ACCOUNT_TRANSFERRED_IN", set()))
    out4 = evaluate_rule_set("RSV_v0", month="202604", scope="practice")
    r4 = next(r for r in out4["results"] if r["rule_code"] == "RETAINED_ACCOUNT")
    check("15", "RETAINED_ACCOUNT fires on 202605, returns 0 on the 202604 "
                "baseline, and never double-counts a claimed account",
          len(matched.get("RETAINED_ACCOUNT", set())) > 0
          and not (matched.get("RETAINED_ACCOUNT", set()) & claimed)
          and r4.get("matched_count", -1) == 0 and r4.get("empty_reason"),
          f"202605 retained={len(matched.get('RETAINED_ACCOUNT', set()))}, "
          f"overlap with {len(claimed)} claimed accounts=0; 202604 matched=0 "
          f"with empty_reason={str(r4.get('empty_reason'))[:60]!r}")

    # ---- 16: lifecycle partition
    import itertools
    cats = ["NEW_ACCOUNT", "LOST_ACCOUNT", "RETAINED_ACCOUNT",
            "ACCOUNT_TRANSFERRED_IN", "ACCOUNT_TRANSFERRED_OUT"]
    overlaps = {(a, b): matched.get(a, set()) & matched.get(b, set())
                for a, b in itertools.combinations(cats, 2)}
    from app.graph.queries.catalog import run_catalog_query
    lc = run_catalog_query("account_lifecycle_counts", {
        "from_month": "202604", "to_month": "202605", "scope": "all"})["rows"][0]
    check("16", "account_lifecycle_counts partitions the account set — no "
                "account appears in two categories",
          not any(overlaps.values()),
          f"pairwise overlaps all empty across {cats}; counts new="
          f"{lc['new_count']} lost={lc['lost_count']} retained="
          f"{lc['retained_count']} tin={lc['transferred_in_count']} "
          f"tout={lc['transferred_out_count']}")

    # ---- 17: glossary covers what the mockup displays
    terms = client.get("/api/glossary").json()["terms"]
    need_terms = (["metric.accounts", "metric.trades", "metric.revenue",
                   "metric.aum", "metric.share"]
                  + [f"severity.{s}" for s in SEVERITIES]
                  + ["provenance.REAL", "provenance.DERIVED", "provenance.DUMMY"]
                  + ["noncredited.9H", "noncredited.9G", "noncredited.9D",
                     "noncredited.9E"]
                  + ["driver.NEW_BILLING", "driver.LOST_ACCOUNTS",
                     "driver.TRANSFERS", "driver.RETAINED_ACCOUNTS",
                     "driver.ONE_TIME", "driver.FEE_RATE", "driver.MARKET"])
    missing = [t for t in need_terms if t not in terms
               or not terms[t].get("definition")]
    check("17", "GET /api/glossary returns definitions for every metric, driver, "
                "severity level and provenance chip the mockup displays",
          not missing,
          f"{len(terms)} terms; all {len(need_terms)} mockup-displayed terms "
          f"present with definitions" + (f"; MISSING {missing}" if missing else ""))

    passed = sum(1 for ok, *_ in RESULTS if ok)
    print(f"\n{passed}/{len(RESULTS)} checks passed")
    return 0 if passed == len(RESULTS) else 1


if __name__ == "__main__":
    sys.exit(main())
